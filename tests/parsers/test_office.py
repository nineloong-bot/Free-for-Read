from io import BytesIO

from docx import Document as DocxDocument
from pptx import Presentation

from free_for_read.parsers.office import PowerPointParser, WordParser
from free_for_read.renderers.markdown import render_markdown


def test_word_parser_maps_headings_and_paragraphs() -> None:
    source = DocxDocument()
    source.add_heading("Chapter One", level=1)
    source.add_paragraph("Once upon a time.")
    buffer = BytesIO()
    source.save(buffer)

    document = WordParser().parse(
        buffer.getvalue(), source_url="https://example.com/book.docx"
    )

    assert document.title == "Chapter One"
    assert render_markdown(document) == "# Chapter One\n\nOnce upon a time."


def test_word_parser_maps_tables() -> None:
    source = DocxDocument()
    table = source.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Role"
    table.cell(1, 0).text = "Ada"
    table.cell(1, 1).text = "Reader"
    buffer = BytesIO()
    source.save(buffer)

    document = WordParser().parse(
        buffer.getvalue(), source_url="https://example.com/table.docx"
    )

    table_node = document.root.children[0]
    assert table_node.type == "table"
    assert table_node.children[0].children[0].text == "Name"
    assert table_node.children[0].children[1].text == "Role"
    assert table_node.children[1].children[0].text == "Ada"
    assert table_node.children[1].children[1].text == "Reader"
    assert render_markdown(document) == (
        "| Name | Role |\n"
        "| --- | --- |\n"
        "| Ada | Reader |"
    )


def test_word_parser_preserves_paragraph_table_order() -> None:
    source = DocxDocument()
    source.add_paragraph("Before table.")
    table = source.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "Name"
    table.cell(0, 1).text = "Role"
    source.add_paragraph("After table.")
    buffer = BytesIO()
    source.save(buffer)

    document = WordParser().parse(
        buffer.getvalue(), source_url="https://example.com/mixed.docx"
    )

    assert [child.type for child in document.root.children] == [
        "paragraph",
        "table",
        "paragraph",
    ]
    assert render_markdown(document) == (
        "Before table.\n\n"
        "| Name | Role |\n"
        "| --- | --- |\n\n"
        "After table."
    )


def test_powerpoint_parser_maps_slides_and_text() -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    slide.shapes.title.text = "Opening"
    slide.placeholders[1].text = "Welcome readers"
    buffer = BytesIO()
    presentation.save(buffer)

    document = PowerPointParser().parse(
        buffer.getvalue(),
        source_url="https://example.com/deck.pptx",
    )

    assert document.title == "Opening"
    slide_node = document.root.children[0]
    assert slide_node.type == "slide"
    assert slide_node.metadata["slide_number"] == 1
    assert slide_node.children[0].type == "heading"
    assert slide_node.children[0].level == 1
    assert slide_node.children[0].text == "Opening"
    assert slide_node.children[1].type == "paragraph"
    assert slide_node.children[1].text == "Welcome readers"
    assert render_markdown(document) == "---\n\n## Slide 1\n\n# Opening\n\nWelcome readers"
