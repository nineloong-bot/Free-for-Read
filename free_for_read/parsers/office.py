from io import BytesIO

from docx import Document as DocxDocument
from docx.table import Table
from docx.text.paragraph import Paragraph
from pptx import Presentation

from free_for_read.core.models import Document, DocumentNode, SourceType


class WordParser:
    source_type = SourceType.WORD

    def parse(self, content: bytes, *, source_url: str) -> Document:
        source = DocxDocument(BytesIO(content))
        root = DocumentNode(type="document")
        for block in source.element.body.iterchildren():
            if block.tag.endswith("}p"):
                node = _paragraph_node(Paragraph(block, source))
                if node is not None:
                    root.children.append(node)
            elif block.tag.endswith("}tbl"):
                root.children.append(_table_node(Table(block, source)))
        return Document(root=root, title=_first_heading(root))


class PowerPointParser:
    source_type = SourceType.POWERPOINT

    def parse(self, content: bytes, *, source_url: str) -> Document:
        presentation = Presentation(BytesIO(content))
        root = DocumentNode(type="document")
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_node = DocumentNode(
                type="slide", metadata={"slide_number": slide_number}
            )
            for shape in slide.shapes:
                if not hasattr(shape, "text"):
                    continue
                text = shape.text.strip()
                if not text:
                    continue
                if shape == slide.shapes.title:
                    slide_node.children.append(
                        DocumentNode(type="heading", text=text, level=1)
                    )
                else:
                    slide_node.children.append(DocumentNode(type="paragraph", text=text))
            root.children.append(slide_node)
        return Document(root=root, title=_first_slide_heading(root))


def _heading_level(style_name: str) -> int:
    try:
        return int(style_name.rsplit(" ", 1)[-1])
    except ValueError:
        return 1


def _paragraph_node(paragraph: Paragraph) -> DocumentNode | None:
    text = paragraph.text.strip()
    if not text:
        return None
    if paragraph.style and paragraph.style.name.startswith("Heading"):
        return DocumentNode(
            type="heading",
            text=text,
            level=_heading_level(paragraph.style.name),
        )
    return DocumentNode(type="paragraph", text=text)


def _table_node(table: Table) -> DocumentNode:
    return DocumentNode(
        type="table",
        children=[
            DocumentNode(
                type="table_row",
                children=[
                    DocumentNode(type="table_cell", text=cell.text.strip())
                    for cell in row.cells
                ],
            )
            for row in table.rows
        ],
    )


def _first_heading(root: DocumentNode) -> str | None:
    for child in root.children:
        if child.type == "heading" and child.text:
            return child.text
    return None


def _first_slide_heading(root: DocumentNode) -> str | None:
    for slide in root.children:
        for child in slide.children:
            if child.type == "heading" and child.text:
                return child.text
    return None
